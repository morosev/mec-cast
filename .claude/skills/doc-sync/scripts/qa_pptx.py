#!/usr/bin/env python3
"""Geometry and content QA for the mec-cast deck.

LibreOffice is not available here, so slides cannot be rendered and eyeballed.
This approximates that pass: estimated text overflow, off-slide shapes, thin
margins, overlaps, plus the deck's content contract (10 slides, Profile B not
mentioned before slide 8).

    python .claude/skills/doc-sync/scripts/qa_pptx.py docs/slides/mec-cast-architecture.pptx

Exit 0 clean, 1 on problems. Needs python-pptx.

Text-fit is estimated from average glyph width, so it is approximate by
construction — it flags real overflow reliably but can be off by a few percent
near the boundary. Treat a marginal hit as "look at that slide", not as truth.
"""
from __future__ import annotations

import math
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("python-pptx not installed:")
    print("  telemetry/python/.venv/bin/pip install python-pptx")
    sys.exit(2)

SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN_MIN = 0.5
CHAR_W = 0.50          # avg glyph width as a fraction of point size
LINE_H = 1.22
EMU = 914400.0

# The deck's content contract.
EXPECTED_SLIDES = 10
PROFILE_B_TERMS = r"\b(webrtc|str0m|sfu|libwebrtc|profile b)\b"
PROFILE_B_FIRST_SLIDE = 8

problems: list[str] = []
info: list[str] = []


def inches(v: int) -> float:
    return v / EMU


def check_geometry(idx: int, slide) -> list[tuple]:
    boxes = []
    for shp in slide.shapes:
        if not (shp.has_text_frame and shp.text_frame.text.strip()):
            continue
        x, y = inches(shp.left), inches(shp.top)
        w, h = inches(shp.width), inches(shp.height)
        text = shp.text_frame.text

        sizes = [
            r.font.size.pt
            for p in shp.text_frame.paragraphs
            for r in p.runs
            if r.font.size
        ]
        fs = max(sizes) if sizes else 12.0

        cpl = max(int(max(w * 72 - 4, 10) / (CHAR_W * fs)), 8)
        lines = sum(max(1, math.ceil(len(par) / cpl)) for par in text.split("\n"))
        needed = lines * fs * LINE_H / 72.0
        if needed > h + 0.04:
            problems.append(
                f"[{idx}] OVERFLOW '{text[:38]}…' needs ~{needed:.2f}in, has {h:.2f}in"
            )

        if (x < MARGIN_MIN - 0.01 or y < MARGIN_MIN - 0.01
                or x + w > SLIDE_W - MARGIN_MIN + 0.01
                or y + h > SLIDE_H - MARGIN_MIN + 0.01):
            problems.append(
                f"[{idx}] MARGIN   '{text[:38]}…' "
                f"({x:.2f},{y:.2f})-({x + w:.2f},{y + h:.2f})"
            )
        if x < -0.01 or y < -0.01 or x + w > SLIDE_W + 0.01 or y + h > SLIDE_H + 0.01:
            problems.append(f"[{idx}] OFF-SLIDE '{text[:38]}…'")

        boxes.append((x, y, w, h, text))
    return boxes


def check_overlaps(idx: int, boxes: list[tuple]) -> None:
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            ax, ay, aw, ah, at = boxes[i]
            bx, by, bw, bh, bt = boxes[j]
            ox = min(ax + aw, bx + bw) - max(ax, bx)
            oy = min(ay + ah, by + bh) - max(ay, by)
            if ox <= 0.05 or oy <= 0.05:
                continue
            # A card's title/body legitimately sits inside the card.
            nested = (
                (ax >= bx - .01 and ay >= by - .01
                 and ax + aw <= bx + bw + .01 and ay + ah <= by + bh + .01)
                or (bx >= ax - .01 and by >= ay - .01
                    and bx + bw <= ax + aw + .01 and by + bh <= ay + ah + .01)
            )
            if not nested:
                problems.append(
                    f"[{idx}] OVERLAP  '{at[:24]}' x '{bt[:24]}' "
                    f"({ox:.2f}x{oy:.2f} in)"
                )


def check_content(texts: list[str]) -> None:
    if len(texts) != EXPECTED_SLIDES:
        problems.append(
            f"deck has {len(texts)} slides, contract says {EXPECTED_SLIDES} — "
            "if this is intentional, the structure change needs to be agreed first"
        )

    for i, t in enumerate(texts[:PROFILE_B_FIRST_SLIDE - 1], 1):
        hits = set(m.lower() for m in re.findall(PROFILE_B_TERMS, t, re.I))
        if hits:
            problems.append(
                f"[{i}] Profile B leaked into slides 1-{PROFILE_B_FIRST_SLIDE - 1}: {sorted(hits)}"
            )

    for i, t in enumerate(texts, 1):
        if re.search(r"(?i)lorem|ipsum|\bTODO\b|\[insert|\bxxx\b", t):
            problems.append(f"[{i}] placeholder text left in the slide")

    counts = [len(t.split()) for t in texts]
    avg = sum(counts) / len(counts) if counts else 0
    for i, c in enumerate(counts, 1):
        flag = "  <-- crowded; consider whether this slide is doing too much" \
            if c > avg * 1.45 else ""
        info.append(f"  slide {i:2d}: {c:3d} words{flag}")


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 2
    deck = Path(sys.argv[1])
    if not deck.exists():
        print(f"not found: {deck}")
        return 2

    prs = Presentation(str(deck))
    texts = []
    for idx, slide in enumerate(prs.slides, 1):
        boxes = check_geometry(idx, slide)
        check_overlaps(idx, boxes)
        texts.append("\n".join(b[4] for b in boxes))

    check_content(texts)

    print(f"{deck.name}: {len(texts)} slides, "
          f"{inches(prs.slide_width):.2f} x {inches(prs.slide_height):.2f} in\n")
    for line in info:
        print(line)
    if problems:
        print()
        for p in problems:
            print(f"  {p}")
        print(f"\n{len(problems)} problem(s)")
        return 1
    print("\nPASS — geometry and content contract both hold")
    return 0


if __name__ == "__main__":
    sys.exit(main())
